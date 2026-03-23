"""
文件解析器工厂模式
支持多种文件格式的差异化解析
"""

from typing import Dict, List, Any, Optional, Union, Protocol, Tuple
from abc import ABC, abstractmethod
from enum import Enum
import asyncio
import base64
import re
import os
from pathlib import Path
from urllib.parse import unquote
import uuid

from app.core.logger import get_logger
from datetime import datetime

logger = get_logger(__name__)

# 常见图片魔术头，用于校验 URL 下载内容是否为图片
_IMAGE_MAGIC = [
    (b"\xff\xd8\xff", "image/jpeg"),
    (b"\x89PNG\r\n\x1a\n", "image/png"),
    (b"GIF87a", "image/gif"),
    (b"GIF89a", "image/gif"),
    (b"RIFF", "image/webp"),  # WEBP 需再检查 8:12 是否为 WEBP
    (b"BM", "image/bmp"),
]


def normalize_text_newlines(text: str) -> str:
    """统一换行为 \\n，便于按 \\n\\n 分段（修复 CRLF 下 split('\\n\\n') 失效）。"""
    if not text:
        return text
    return text.replace("\r\n", "\n").replace("\r", "\n")


def _is_image_bytes(data: bytes) -> bool:
    """根据魔术头判断是否为常见图片格式。"""
    if len(data) < 12:
        return data.startswith(b"\xff\xd8\xff") or data.startswith(b"\x89PNG") or data[:6] in (b"GIF87a", b"GIF89a") or data.startswith(b"BM")
    if data[:4] == b"RIFF" and data[8:12] == b"WEBP":
        return True
    for magic, _ in _IMAGE_MAGIC:
        if data.startswith(magic):
            return True
    return False


async def _fetch_image_from_url(url: str, timeout: int, max_size: int) -> Optional[bytes]:
    """异步下载 URL，校验为图片且不超过 max_size，返回 bytes 或 None。"""
    try:
        import httpx
        async with httpx.AsyncClient(timeout=timeout, follow_redirects=True) as client:
            r = await client.get(url)
            r.raise_for_status()
            content = r.content
            if len(content) > max_size:
                logger.warning("Markdown 图片 URL 超过大小限制，已跳过: {} (size={})", url[:80], len(content))
                return None
            ct = (r.headers.get("content-type") or "").split(";")[0].strip().lower()
            if not _is_image_bytes(content) and not ct.startswith("image/"):
                logger.debug("Markdown 图片 URL 非图片类型，已跳过: {} (content-type={})", url[:80], ct)
                return None
            return content
    except Exception as e:
        logger.debug("Markdown 图片 URL 下载失败: {} ({})", url[:80], e)
        return None


def _resolve_relative_path(file_path: str, ref_path: str) -> str:
    """将相对 ref_path 相对 file_path 所在目录解析为规范路径（与 asset_map 的 key 一致）。"""
    base = Path(file_path).parent
    # 去掉 ref 中的 leading ./
    ref = ref_path.strip().lstrip("./")
    resolved = (base / ref).as_posix()
    return resolved


def _read_local_image_if_allowed(
    path: str,
    allowed_base_paths: List[str],
    max_size: int,
) -> Optional[bytes]:
    """
    当 path 为本地绝对路径（/path 或 file:///path）且位于 allowed_base_paths 之下时，读取文件并返回字节。
    """
    if not allowed_base_paths:
        return None
    path = path.strip()
    if path.lower().startswith("file://"):
        path = unquote(path[7:].lstrip("/"))
        if path and not path.startswith("/"):
            path = os.path.abspath("/" + path) if os.name != "nt" else path
    if not path or (not path.startswith("/") and len(path) < 2):
        return None
    if os.name == "nt" and path.startswith("/") and not path.startswith("//"):
        path = path.lstrip("/")
    try:
        p = Path(path).resolve()
        if not p.is_file():
            return None
        real_str = str(p)
        allowed = False
        for base in allowed_base_paths:
            base_p = Path(base).resolve()
            base_str = str(base_p)
            if real_str == base_str or real_str.startswith(base_str + os.sep):
                allowed = True
                break
        if not allowed:
            return None
        data = p.read_bytes()
        if len(data) > max_size:
            return None
        if not _is_image_bytes(data):
            return None
        return data
    except Exception as e:
        logger.debug("读取本地图片失败 {}: {}", path[:80], e)
        return None


async def _extract_all_images_from_markdown(
    content: str,
    file_path: str = "",
    asset_map: Optional[Dict[str, bytes]] = None,
    fetch_urls: bool = True,
    timeout: int = 10,
    max_size: int = 5 * 1024 * 1024,
    allowed_local_base_paths: Optional[List[str]] = None,
) -> List[Dict[str, Any]]:
    """
    从 Markdown 中按出现顺序提取所有图片：data: base64、http(s) 链接、相对路径（需提供 asset_map）。
    返回与 PDF/DOCX/PPTX 兼容的 extracted_images 项列表。
    """
    refs_ordered = _extract_image_refs_from_markdown_text(content)
    if refs_ordered and allowed_local_base_paths:
        logger.info(
            "Markdown 发现 {} 处图片引用，本地路径白名单已配置（{} 项），将尝试读取本地图片",
            len(refs_ordered),
            len(allowed_local_base_paths),
        )
    extracted: List[Dict[str, Any]] = []
    for idx, (ref_string, path) in enumerate(refs_ordered):
        path = path.strip()
        image_bytes: Optional[bytes] = None
        source: str = "markdown"
        if path.lower().startswith("data:"):
            if ";base64," in path:
                try:
                    b64 = path.split(";base64,", 1)[1]
                    image_bytes = base64.b64decode(b64, validate=True)
                    source = "markdown-inline-base64"
                except Exception as e:
                    logger.debug("跳过无效 base64 图片: {}", e)
                    continue
        elif path.startswith("http://") or path.startswith("https://"):
            if fetch_urls:
                image_bytes = await _fetch_image_from_url(path, timeout, max_size)
                source = "markdown-url"
            else:
                continue
        elif path.startswith("/") or path.lower().startswith("file://"):
            # 本地绝对路径：仅在配置白名单内时读取
            if allowed_local_base_paths:
                image_bytes = _read_local_image_if_allowed(path, allowed_local_base_paths, max_size)
                if image_bytes is not None:
                    source = "markdown-local-path"
                    logger.info("Markdown 已从本地路径读取图片: {}", path[:80] + ("..." if len(path) > 80 else ""))
                else:
                    logger.warning("Markdown 本地图片未读取（路径不在白名单下或文件不存在）: {}", path[:80] + ("..." if len(path) > 80 else ""))
        else:
            # 相对路径：仅当提供 asset_map 时解析
            if asset_map:
                key = _resolve_relative_path(file_path, path)
                image_bytes = asset_map.get(key)
                if image_bytes is None:
                    # 尝试无前导 ./ 的 key
                    key_alt = (Path(file_path).parent / path.strip().lstrip("./")).as_posix()
                    image_bytes = asset_map.get(key_alt)
                if image_bytes is not None:
                    source = "markdown-relative"
        if image_bytes and len(image_bytes) > 0:
            extracted.append({
                "page": 1,
                "image_index": len(extracted),
                "image_bytes": image_bytes,
                "image_path": path[:200] if not path.startswith("data:") else "",
                "markdown_ref": ref_string,
                "metadata": {"extracted_from": source},
            })
    return extracted


def _extract_image_refs_from_markdown_text(text: str) -> List[Tuple[str, str]]:
    """
    从 markdown/HTML 文本中按出现顺序提取图片引用，返回 [(完整引用串, 路径), ...]。
    支持 Markdown ![](path) 与 HTML <img src="path" ...>。
    """
    # (完整引用串, 路径, 起始位置)
    with_pos: List[Tuple[str, str, int]] = []
    for m in re.finditer(r"!\[[^\]]*\]\s*\(\s*([^)]+)\s*\)", text):
        with_pos.append((m.group(0), m.group(1).strip(), m.start()))
    for m in re.finditer(r'<img[^>]+src=["\']([^"\']+)["\'][^>]*>', text):
        with_pos.append((m.group(0), m.group(1).strip(), m.start()))
    with_pos.sort(key=lambda x: x[2])
    return [(r[0], r[1]) for r in with_pos]

class FileType(Enum):
    """文件类型枚举"""
    PDF = "pdf"
    DOCX = "docx"
    DOC = "doc"
    PPTX = "pptx"
    TXT = "txt"
    MD = "md"
    IMAGE = "image"
    AUDIO = "audio"      # mp3, wav, m4a, flac, aac, ogg等
    VIDEO = "video"      # mp4, avi, mov, mkv, webm等
    UNKNOWN = "unknown"

class DocumentParser(ABC):
    """文档解析器基类"""
    
    @abstractmethod
    async def parse(self, file_content: bytes, file_path: str, **kwargs: Any) -> Dict[str, Any]:
        """
        解析文件内容
        
        Args:
            file_content: 文件二进制内容
            file_path: 文件路径
            **kwargs: 扩展参数（如 asset_map 供 Markdown 相对路径图片解析）
            
        Returns:
            解析结果字典
        """
        pass
    
    @abstractmethod
    def supports_file_type(self) -> FileType:
        """支持的文件类型"""
        pass

class PDFParser(DocumentParser):
    """PDF文档解析器"""
    
    def __init__(self):
        self.supported_formats = ["pdf"]
    
    def supports_file_type(self) -> FileType:
        return FileType.PDF
    
    async def parse(self, file_content: bytes, file_path: str, **kwargs: Any) -> Dict[str, Any]:
        """解析PDF：优先 MinerU API，失败则本地 MinerU2.5，再 PaddleOCR-VL-1.5，最后 PyMuPDF"""
        from app.core.config import settings

        # 1. 优先 MinerU API（需 MINERU_TOKEN）
        token = getattr(settings, "mineru_token", None)
        if token and getattr(settings, "mineru_pdf_enabled", True):
            try:
                from .mineru_client import parse_pdf_via_api
                loop = asyncio.get_event_loop()
                result = await loop.run_in_executor(
                    None,
                    lambda: parse_pdf_via_api(file_content, file_path, token),
                )
                if result is not None:
                    logger.info("使用 MinerU API 解析 PDF")
                    return result
            except Exception as e:
                logger.warning("MinerU API 解析失败，尝试本地模型: {}", e)

        # 2. 备选 本地 MinerU2.5
        try:
            from .mineru_client import get_mineru_client, parse_pdf as mineru_parse_pdf
            if getattr(settings, "mineru_pdf_enabled", True):
                client = get_mineru_client()
                if client:
                    logger.info("使用 MinerU2.5 解析 PDF")
                    loop = asyncio.get_event_loop()
                    return await loop.run_in_executor(
                        None, lambda: mineru_parse_pdf(file_content)
                    )
        except Exception as e:
            logger.warning("MinerU 本地解析失败，尝试 PaddleOCR: {}", e)

        # 3. 备选 PaddleOCR-VL-1.5
        try:
            from .paddleocr_client import get_paddleocr_client

            ocr_client = get_paddleocr_client()
            if ocr_client:
                logger.info("使用 PaddleOCR-VL-1.5 解析 PDF")
                return await self._parse_with_paddleocr(file_content, file_path)
        except Exception as e:
            logger.warning("PaddleOCR 解析失败，回退到 PyMuPDF: {}", e)

        # 4. 回退 PyMuPDF
        return await self._parse_with_pymupdf(file_content, file_path)
    
    async def _parse_with_paddleocr(self, file_content: bytes, file_path: str) -> Dict[str, Any]:
        """使用 PaddleOCR-VL-1.5 解析 PDF"""
        from .paddleocr_client import get_paddleocr_client
        
        client = get_paddleocr_client()
        if not client:
            raise ValueError("PaddleOCR 客户端不可用")
        
        from app.core.config import settings
        max_pixels = getattr(settings, "paddleocr_max_pixels", None)
        result_data = client.parse_pdf(
            file_content=file_content,
            file_type=0,  # PDF
            use_doc_orientation_classify=False,
            use_doc_unwarping=False,
            use_chart_recognition=False,
            max_pixels=max_pixels,
        )
        
        layout_results = result_data.get("layoutParsingResults", [])
        total_pages = len(layout_results)
        
        # 提取每页的 markdown 和图片信息
        pages_content = []
        extracted_images = []
        markdown_parts = []
        
        for page_index, page_result in enumerate(layout_results):
            page_num = page_index + 1
            markdown_data = page_result.get("markdown", {})
            markdown_text = markdown_data.get("text", "")
            markdown_images = markdown_data.get("images", {})
            # 文档标题/说明：API 若有 image_captions 或 imageCaptions 则填入
            image_captions = markdown_data.get("image_captions") or markdown_data.get("imageCaptions") or page_result.get("image_captions") or page_result.get("imageCaptions") or {}

            # 按 markdown 中图片引用出现顺序解析，保证与 full_markdown 一一对应，便于 VLM 图注插回
            refs_ordered = _extract_image_refs_from_markdown_text(markdown_text)
            page_images = []
            for img_index, (ref_string, path_in_text) in enumerate(refs_ordered):
                # 用路径匹配 markdown_images 的 key（支持路径、反斜杠、basename）
                path_n = path_in_text.replace("\\", "/")
                img_url = markdown_images.get(path_in_text) or markdown_images.get(path_n)
                if not img_url and path_in_text:
                    for k, v in markdown_images.items():
                        if Path(k).name == Path(path_in_text).name:
                            img_url = v
                            break
                if not img_url:
                    logger.debug("PaddleOCR 页 {} 中引用路径未在 images 中找到: {}", page_num, path_in_text)
                    continue
                try:
                    image_bytes = client.download_image(img_url)
                    img_path = path_in_text
                    doc_caption = None
                    if isinstance(image_captions, dict):
                        doc_caption = image_captions.get(path_in_text) or image_captions.get(path_n) or image_captions.get(Path(path_in_text).name)
                    if isinstance(image_captions, list) and img_index < len(image_captions):
                        cap = image_captions[img_index]
                        doc_caption = cap if isinstance(cap, str) else (cap.get("text") or cap.get("caption") if isinstance(cap, dict) else None)

                    image_info = {
                        "page": page_num,
                        "image_index": len(page_images),
                        "image_bytes": image_bytes,
                        "image_path": Path(img_path).name or f"page{page_num}_img{len(page_images)}.jpg",
                        "markdown_ref": ref_string,
                        "metadata": {
                            "extracted_from": "paddleocr",
                            "page": page_num,
                            "document_caption": doc_caption,
                        },
                    }
                    extracted_images.append(image_info)
                    page_images.append(image_info)
                except Exception as e:
                    logger.warning("下载第 {} 页图片失败 (ref={}): {}", page_num, path_in_text[:50], e)
            
            # 构建页面内容
            pages_content.append({
                "page": page_num,
                "markdown": markdown_text,
                "text": markdown_text,
                "images": page_images,
                "metadata": {"parser": "paddleocr-vl-1.5"},
            })
            if markdown_text.strip():
                if markdown_parts:
                    markdown_parts.append(f"\n\n---\n\n## 第 {page_num} 页\n\n")
                markdown_parts.append(markdown_text)
        
        # 拼接完整的 markdown
        full_markdown = "\n\n".join(markdown_parts)
        
        return {
            "file_type": "pdf",
            "markdown": full_markdown,
            "pages": pages_content,
            "extracted_images": extracted_images,
            "total_pages": total_pages,
            "metadata": {
                "parser": "paddleocr-vl-1.5",
                "extracted_at": datetime.utcnow().isoformat() + "Z"
            }
        }
    
    async def _parse_with_pymupdf(self, file_content: bytes, file_path: str) -> Dict[str, Any]:
        """使用 PyMuPDF 解析 PDF（原有方法）"""
        try:
            import fitz  # PyMuPDF
            
            # 保存临时文件
            temp_file = f"/tmp/{uuid.uuid4()}.pdf"
            with open(temp_file, "wb") as f:
                f.write(file_content)
            
            # 使用PyMuPDF解析
            doc = fitz.open(temp_file)
            total_pages = len(doc)
            pages_content = []
            tables = []
            
            for page_num in range(total_pages):
                page = doc.load_page(page_num)
                
                # 提取文本
                text = page.get_text()
                
                # 尝试提取表格（使用pdfplumber，如果可用）
                try:
                    import pdfplumber  # type: ignore
                    with pdfplumber.open(temp_file) as pdf:
                        pdf_page = pdf.pages[page_num]
                        tables_on_page = pdf_page.find_tables()
                        for table in tables_on_page:
                            table_data = table.extract()
                            if table_data:
                                tables.append(table_data)
                except (ImportError, Exception):
                    # pdfplumber不可用或提取失败，跳过表格提取
                    pass
                
                pages_content.append({
                    "page": page_num + 1,
                    "text": text,
                    "metadata": {
                        "width": page.rect.width,
                        "height": page.rect.height
                    }
                })
            
            doc.close()
            
            # 清理临时文件
            Path(temp_file).unlink(missing_ok=True)
            
            return {
                "file_type": "pdf",
                "pages": pages_content,
                "tables": tables,
                "total_pages": total_pages,
                "metadata": {
                    "parser": "pymupdf",
                    "extracted_at": datetime.utcnow().isoformat() + "Z"
                }
            }
            
        except Exception as e:
            logger.error(f"PDF解析失败: {str(e)}")
            raise
    
    async def parse_with_pymupdf4llm(self, file_content: bytes, file_path: str) -> Dict[str, Any]:
        """使用pymupdf4llm解析（如果可用）"""
        try:
            import pymupdf4llm
            
            # 保存临时文件
            temp_file = f"/tmp/{uuid.uuid4()}.pdf"
            with open(temp_file, "wb") as f:
                f.write(file_content)
            
            # 使用pymupdf4llm解析
            md_text = pymupdf4llm.to_markdown(temp_file)
            
            # 清理临时文件
            Path(temp_file).unlink(missing_ok=True)
            
            return {
                "file_type": "pdf",
                "markdown": md_text,
                "pages": [{"page": 1, "text": md_text}],
                "metadata": {
                    "parser": "pymupdf4llm",
                    "extracted_at": "2024-01-01T00:00:00Z"
                }
            }
            
        except ImportError:
            logger.warning("pymupdf4llm不可用，使用PyMuPDF")
            return await self.parse(file_content, file_path)
        except Exception as e:
            logger.error(f"pymupdf4llm解析失败: {str(e)}")
            return await self.parse(file_content, file_path)

class DocxParser(DocumentParser):
    """DOCX文档解析器：优先 MinerU API，失败则本地 MinerU2.5（需 LibreOffice），最后 python-docx"""

    def supports_file_type(self) -> FileType:
        return FileType.DOCX

    async def parse(self, file_content: bytes, file_path: str, **kwargs: Any) -> Dict[str, Any]:
        from app.core.config import settings

        # 1. 优先 MinerU API
        token = getattr(settings, "mineru_token", None)
        if token and getattr(settings, "mineru_pdf_enabled", True):
            try:
                from .mineru_client import parse_docx_via_api
                loop = asyncio.get_event_loop()
                result = await loop.run_in_executor(
                    None,
                    lambda: parse_docx_via_api(file_content, file_path, token),
                )
                if result is not None:
                    logger.info("使用 MinerU API 解析 Word(docx)")
                    return result
            except Exception as e:
                logger.warning("MinerU API 解析 docx 失败，尝试本地 MinerU: {}", e)

        # 2. 备选 本地 MinerU2.5（需 LibreOffice 将 docx 转 PDF 再按页 MinerU 提取）
        try:
            from .mineru_client import get_mineru_client, parse_docx as mineru_parse_docx
            if getattr(settings, "mineru_pdf_enabled", True):
                client = get_mineru_client()
                if client:
                    logger.info("使用 MinerU2.5 解析 Word(docx)")
                    loop = asyncio.get_event_loop()
                    return await loop.run_in_executor(
                        None, lambda: mineru_parse_docx(file_content)
                    )
        except Exception as e:
            logger.warning("MinerU 本地解析 docx 失败，回退 python-docx: {}", e)

        # 3. 回退 python-docx
        return await self._parse_with_docx(file_content, file_path)

    async def _parse_with_docx(self, file_content: bytes, file_path: str) -> Dict[str, Any]:
        """使用 python-docx 解析 DOCX（原有逻辑）"""
        try:
            from docx import Document
            import io

            doc = Document(io.BytesIO(file_content))
            paragraphs = []
            tables = []

            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append({
                        "text": para.text.strip(),
                        "style": para.style.name if para.style else "Normal"
                    })
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip())
                    table_data.append(row_data)
                tables.append(table_data)

            return {
                "file_type": "docx",
                "paragraphs": paragraphs,
                "tables": tables,
                "metadata": {
                    "total_paragraphs": len(paragraphs),
                    "total_tables": len(tables),
                    "extracted_at": datetime.utcnow().isoformat() + "Z"
                }
            }
        except Exception as e:
            logger.error("DOCX 解析失败: {}", e)
            raise


class PptxParser(DocumentParser):
    """PPTX文档解析器：优先 MinerU API，失败则本地 MinerU2.5（需 LibreOffice），最后 python-pptx"""

    def supports_file_type(self) -> FileType:
        return FileType.PPTX

    async def parse(self, file_content: bytes, file_path: str, **kwargs: Any) -> Dict[str, Any]:
        from app.core.config import settings

        # 1. 优先 MinerU API
        token = getattr(settings, "mineru_token", None)
        if token and getattr(settings, "mineru_pdf_enabled", True):
            try:
                from .mineru_client import parse_pptx_via_api
                loop = asyncio.get_event_loop()
                result = await loop.run_in_executor(
                    None,
                    lambda: parse_pptx_via_api(file_content, file_path, token),
                )
                if result is not None:
                    logger.info("使用 MinerU API 解析 PPTX")
                    return result
            except Exception as e:
                logger.warning("MinerU API 解析 pptx 失败，尝试本地 MinerU: {}", e)

        # 2. 备选 本地 MinerU2.5（需 LibreOffice 将 pptx 转 PDF 再按页 MinerU 提取）
        try:
            from .mineru_client import get_mineru_client, parse_pptx as mineru_parse_pptx
            if getattr(settings, "mineru_pdf_enabled", True):
                client = get_mineru_client()
                if client:
                    logger.info("使用 MinerU2.5 解析 PPTX")
                    loop = asyncio.get_event_loop()
                    return await loop.run_in_executor(
                        None, lambda: mineru_parse_pptx(file_content)
                    )
        except Exception as e:
            logger.warning("MinerU 本地解析 pptx 失败，回退 python-pptx: {}", e)

        # 3. 回退 python-pptx（仅文本）
        return await self._parse_with_pptx(file_content, file_path)

    async def _parse_with_pptx(self, file_content: bytes, file_path: str) -> Dict[str, Any]:
        """使用 python-pptx 解析 PPTX（仅文本，无布局/图片）"""
        try:
            import io
            from pptx import Presentation

            prs = Presentation(io.BytesIO(file_content))
            paragraphs = []
            for slide_idx, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    text = getattr(shape, "text", "")
                    if text and text.strip():
                        paragraphs.append({
                            "text": text.strip(),
                            "slide_index": slide_idx + 1,
                        })
            return {
                "file_type": "pptx",
                "paragraphs": paragraphs,
                "metadata": {
                    "total_paragraphs": len(paragraphs),
                    "total_slides": len(prs.slides),
                    "extracted_at": datetime.utcnow().isoformat() + "Z"
                }
            }
        except Exception as e:
            logger.error("PPTX 解析失败: {}", e)
            raise


class TextParser(DocumentParser):
    """文本文件解析器"""
    
    def supports_file_type(self) -> FileType:
        return FileType.TXT
    
    async def parse(self, file_content: bytes, file_path: str, **kwargs: Any) -> Dict[str, Any]:
        """解析文本文件"""
        try:
            # 检测编码
            content = normalize_text_newlines(file_content.decode('utf-8', errors='ignore'))
            
            # 简单分段落（须先归一化换行，否则 CRLF 文件无法按 \n\n 切开）
            paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]
            
            return {
                "file_type": "txt",
                "content": content,
                "paragraphs": [{"text": para} for para in paragraphs],
                "metadata": {
                    "encoding": "utf-8",
                    "total_characters": len(content),
                    "total_paragraphs": len(paragraphs),
                    "extracted_at": "2024-01-01T00:00:00Z"
                }
            }
            
        except Exception as e:
            logger.error(f"文本文件解析失败: {str(e)}")
            raise

class MarkdownParser(DocumentParser):
    """Markdown文档解析器"""
    
    def supports_file_type(self) -> FileType:
        return FileType.MD
    
    async def parse(self, file_content: bytes, file_path: str, **kwargs: Any) -> Dict[str, Any]:
        """解析 Markdown 文件；支持通过 kwargs.asset_map 传入相对路径图片字节（如文件夹导入）。"""
        try:
            content = normalize_text_newlines(file_content.decode('utf-8', errors='ignore'))
            
            # 使用markdown库解析
            import markdown
            
            html_content = markdown.markdown(content)
            
            # 提取结构化信息
            lines = content.split('\n')
            headers = []
            code_blocks = []
            current_code_block = []
            in_code_block = False
            
            for line in lines:
                # 检测标题
                if line.startswith('#'):
                    headers.append({
                        "level": len(line) - len(line.lstrip('#')),
                        "text": line.lstrip('#').strip()
                    })
                
                # 检测代码块
                if line.strip().startswith('```'):
                    if in_code_block:
                        if current_code_block:
                            code_blocks.append('\n'.join(current_code_block))
                            current_code_block = []
                    in_code_block = not in_code_block
                elif in_code_block:
                    current_code_block.append(line)
            
            if current_code_block:
                code_blocks.append('\n'.join(current_code_block))
            
            # 生成智能段落：将标题和紧随的内容绑定在一起
            paragraphs = self._build_smart_paragraphs(content, headers)
            
            # 提取所有图片（base64、http(s)、相对路径 asset_map、本地绝对路径白名单），与 PDF/DOCX/PPTX 一致走 VLM/CLIP 流水线
            from app.core.config import settings
            fetch_urls = getattr(settings, "markdown_fetch_image_urls", True)
            timeout = getattr(settings, "markdown_image_url_timeout", 10)
            max_size = getattr(settings, "markdown_image_url_max_size", 5 * 1024 * 1024)
            asset_map = (kwargs.get("asset_map") or None) if kwargs else None
            allowed_local = getattr(settings, "markdown_local_image_allowed_base_paths", None) or []
            extracted_images = await _extract_all_images_from_markdown(
                content,
                file_path=file_path,
                asset_map=asset_map,
                fetch_urls=fetch_urls,
                timeout=timeout,
                max_size=max_size,
                allowed_local_base_paths=allowed_local if allowed_local else None,
            )
            if extracted_images:
                logger.info("从 Markdown 中提取 {} 张图片（含链接/内联）", len(extracted_images))
            
            result: Dict[str, Any] = {
                "file_type": "md",
                "content": content,
                "html_content": html_content,
                "headers": headers,
                "code_blocks": code_blocks,
                "paragraphs": paragraphs,  # 添加智能段落
                "markdown": content,  # 供 service 分块与图片上下文使用
                "metadata": {
                    "parser": "markdown",
                    "total_headers": len(headers),
                    "total_code_blocks": len(code_blocks),
                    "total_paragraphs": len(paragraphs),
                    "extracted_at": "2024-01-01T00:00:00Z"
                }
            }
            if extracted_images:
                result["extracted_images"] = extracted_images
            return result
            
        except Exception as e:
            logger.error(f"Markdown解析失败: {str(e)}")
            raise
    
    def _build_smart_paragraphs(self, content: str, headers: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """构建智能段落：将标题和紧随的内容绑定在一起
        
        策略：
        1. 标题总是和紧随其后的内容绑定在一起
        2. 遇到新的标题时，结束当前段落并开始新段落
        3. 代码块作为独立段落处理
        4. 空行用于分隔段落，但如果标题后只有空行，继续等待内容
        
        Args:
            content: Markdown内容
            headers: 标题列表（用于参考，实际从内容中解析）
            
        Returns:
            段落列表，每个段落包含标题（如果有）和内容
        """
        lines = content.split('\n')
        # 无任何 ATX 标题时，单空行即分段（公文/中文编号小节等无 # 结构）
        has_atx_headers = any(l.startswith('#') for l in lines)
        paragraphs = []
        current_paragraph = []
        current_header = None
        in_code_block = False
        code_block_language = None
        
        i = 0
        while i < len(lines):
            line = lines[i]
            stripped_line = line.strip()
            
            # 检测代码块
            if stripped_line.startswith('```'):
                if in_code_block:
                    # 结束代码块
                    current_paragraph.append(line)
                    # 保存代码块段落
                    paragraph_text = '\n'.join(current_paragraph).strip()
                    if paragraph_text:
                        paragraphs.append({
                            "text": paragraph_text,
                            "header": None,  # 代码块不绑定标题
                            "has_code": True
                        })
                    current_paragraph = []
                    in_code_block = False
                    code_block_language = None
                else:
                    # 开始代码块
                    # 如果当前有段落，先保存（不包含代码块开始标记）
                    if current_paragraph:
                        paragraph_text = '\n'.join(current_paragraph).strip()
                        if paragraph_text:
                            paragraphs.append({
                                "text": paragraph_text,
                                "header": current_header,
                                "has_code": False
                            })
                        current_paragraph = []
                        current_header = None
                    # 开始新的代码块段落
                    in_code_block = True
                    code_block_language = stripped_line[3:].strip()  # 提取语言标识
                    current_paragraph = [line]
                i += 1
                continue
            
            if in_code_block:
                # 在代码块中，直接添加
                current_paragraph.append(line)
                i += 1
                continue
            
            # 检测标题
            if line.startswith('#'):
                # 如果当前有段落，先保存
                if current_paragraph:
                    paragraph_text = '\n'.join(current_paragraph).strip()
                    if paragraph_text:
                        paragraphs.append({
                            "text": paragraph_text,
                            "header": current_header,
                            "has_code": False
                        })
                    current_paragraph = []
                
                # 提取标题信息
                header_level = len(line) - len(line.lstrip('#'))
                header_text = line.lstrip('#').strip()
                current_header = {
                    "level": header_level,
                    "text": header_text
                }
                
                # 标题本身作为段落的第一行
                current_paragraph.append(line)
                i += 1
                continue
            
            # 空行处理
            if not stripped_line:
                if not has_atx_headers:
                    if current_paragraph:
                        paragraph_text = '\n'.join(current_paragraph).strip()
                        if paragraph_text:
                            paragraphs.append({
                                "text": paragraph_text,
                                "header": current_header,
                                "has_code": False
                            })
                        current_paragraph = []
                        current_header = None
                    i += 1
                    continue
                # 检查下一行是否是标题
                if i + 1 < len(lines):
                    next_line = lines[i + 1]
                    if next_line.startswith('#'):
                        # 下一行是标题，结束当前段落（不包含这个空行）
                        if current_paragraph:
                            paragraph_text = '\n'.join(current_paragraph).strip()
                            if paragraph_text:
                                paragraphs.append({
                                    "text": paragraph_text,
                                    "header": current_header,
                                    "has_code": False
                                })
                            current_paragraph = []
                            current_header = None
                        i += 1
                        continue
                    else:
                        # 下一行不是标题，空行可能是段落分隔符
                        # 但如果当前段落只有标题，继续等待内容
                        non_empty_lines = [l for l in current_paragraph if l.strip()]
                        if len(non_empty_lines) == 1 and non_empty_lines[0].startswith('#'):
                            # 只有标题，添加空行并继续
                            current_paragraph.append(line)
                            i += 1
                            continue
                        else:
                            # 有内容，空行作为段落分隔符
                            # 但先不结束，看看后面是否还有内容
                            # 如果连续两个空行，结束段落
                            if i + 1 < len(lines) and not lines[i + 1].strip():
                                # 连续两个空行，结束段落
                                if current_paragraph:
                                    paragraph_text = '\n'.join(current_paragraph).strip()
                                    if paragraph_text:
                                        paragraphs.append({
                                            "text": paragraph_text,
                                            "header": current_header,
                                            "has_code": False
                                        })
                                    current_paragraph = []
                                    current_header = None
                                i += 2  # 跳过两个空行
                                continue
                            else:
                                # 单个空行，添加到当前段落
                                current_paragraph.append(line)
                                i += 1
                                continue
                else:
                    # 最后一行是空行，结束当前段落
                    if current_paragraph:
                        paragraph_text = '\n'.join(current_paragraph).strip()
                        if paragraph_text:
                            paragraphs.append({
                                "text": paragraph_text,
                                "header": current_header,
                                "has_code": False
                            })
                        current_paragraph = []
                        current_header = None
                    i += 1
                    continue
            
            # 普通内容行
            current_paragraph.append(line)
            i += 1
        
        # 处理最后一个段落
        if current_paragraph:
            paragraph_text = '\n'.join(current_paragraph).strip()
            if paragraph_text:
                paragraphs.append({
                    "text": paragraph_text,
                    "header": current_header,
                    "has_code": in_code_block
                })
        
        return paragraphs

class ImageParser(DocumentParser):
    """图片解析器（用于生成描述）"""
    
    def supports_file_type(self) -> FileType:
        return FileType.IMAGE
    
    async def parse(self, file_content: bytes, file_path: str, **kwargs: Any) -> Dict[str, Any]:
        """解析图片文件"""
        try:
            import base64
            import io
            from PIL import Image
            
            # 获取图片信息
            image = Image.open(io.BytesIO(file_content))
            
            # 转换为base64用于VLM处理
            base64_content = base64.b64encode(file_content).decode('utf-8')
            
            return {
                "file_type": "image",
                "width": image.width,
                "height": image.height,
                "format": image.format,
                "mode": image.mode,
                "base64_content": base64_content,
                "metadata": {
                    "size_bytes": len(file_content),
                    "aspect_ratio": round(image.width / image.height, 2) if image.height > 0 else 0,
                    "extracted_at": "2024-01-01T00:00:00Z"
                }
            }
            
        except Exception as e:
            logger.error(f"图片解析失败: {str(e)}")
            raise

class AudioParser(DocumentParser):
    """音频解析器"""
    
    def supports_file_type(self) -> FileType:
        return FileType.AUDIO
    
    async def parse(self, file_content: bytes, file_path: str, **kwargs: Any) -> Dict[str, Any]:
        """解析音频文件，提取元数据"""
        try:
            import librosa
            import soundfile as sf
            from io import BytesIO
            
            # 使用librosa读取音频元数据
            audio_io = BytesIO(file_content)
            audio_io.seek(0)
            
            # 尝试使用soundfile读取元数据（更快，不需要解码音频）
            try:
                with sf.SoundFile(audio_io) as audio_file:
                    duration = len(audio_file) / audio_file.samplerate
                    sample_rate = audio_file.samplerate
                    channels = audio_file.channels
                    format_name = audio_file.format
                    subtype = audio_file.subtype
            except Exception:
                # 如果soundfile失败，使用librosa（需要解码，较慢）
                audio_io.seek(0)
                y, sample_rate = librosa.load(audio_io, sr=None, duration=0.1)  # 只加载一小段来获取元数据
                duration = librosa.get_duration(y=y, sr=sample_rate)
                channels = 1  # librosa默认转换为单声道
                format_name = Path(file_path).suffix.lower().lstrip('.')
                subtype = None
            
            # 估算比特率（如果无法直接获取）
            file_size = len(file_content)
            bitrate = int((file_size * 8) / duration) if duration > 0 else 0  # bps
            bitrate_kbps = bitrate // 1000  # 转换为kbps
            
            return {
                "file_type": "audio",
                "duration": float(duration),
                "sample_rate": int(sample_rate),
                "channels": int(channels),
                "format": format_name or Path(file_path).suffix.lower().lstrip('.'),
                "bitrate": bitrate_kbps,
                "file_size": file_size,
                "metadata": {
                    "subtype": subtype,
                    "extracted_at": datetime.now().isoformat()
                }
            }
            
        except Exception as e:
            logger.error(f"音频解析失败: {str(e)}")
            # 如果解析失败，返回基本信息
            return {
                "file_type": "audio",
                "duration": 0.0,
                "sample_rate": 44100,  # 默认值
                "channels": 2,  # 默认值
                "format": Path(file_path).suffix.lower().lstrip('.'),
                "bitrate": 0,
                "file_size": len(file_content),
                "metadata": {
                    "extracted_at": datetime.now().isoformat(),
                    "parse_error": str(e)
                }
            }

class VideoParser(DocumentParser):
    """视频解析器"""
    
    def supports_file_type(self) -> FileType:
        return FileType.VIDEO
    
    async def parse(self, file_content: bytes, file_path: str, **kwargs: Any) -> Dict[str, Any]:
        """解析视频文件，提取元数据"""
        try:
            import cv2
            import numpy as np
            from io import BytesIO
            
            # 将bytes转换为numpy数组
            video_bytes = np.frombuffer(file_content, dtype=np.uint8)
            
            # 使用OpenCV读取视频信息
            # 注意：cv2.VideoCapture需要文件路径或URL，不能直接读取bytes
            # 我们需要先保存到临时文件或使用其他方法
            import tempfile
            import os
            
            with tempfile.NamedTemporaryFile(delete=False, suffix=Path(file_path).suffix) as tmp_file:
                tmp_file.write(file_content)
                tmp_path = tmp_file.name
            
            try:
                cap = cv2.VideoCapture(tmp_path)
                
                if not cap.isOpened():
                    raise ValueError("无法打开视频文件")
                
                # 获取视频属性
                fps = cap.get(cv2.CAP_PROP_FPS)
                frame_count = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
                width = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
                height = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
                duration = frame_count / fps if fps > 0 else 0
                
                # 检查是否有音频轨道
                has_audio = False
                try:
                    # 尝试读取第一帧来检测音频
                    # OpenCV无法直接检测音频，我们假设如果视频文件较大可能包含音频
                    file_size = len(file_content)
                    has_audio = file_size > (width * height * frame_count * 3)  # 粗略估算
                except:
                    pass
                
                # 获取编码格式（如果可用）
                fourcc = int(cap.get(cv2.CAP_PROP_FOURCC))
                codec = "".join([chr((fourcc >> 8 * i) & 0xFF) for i in range(4)])
                
                cap.release()
                
                return {
                    "file_type": "video",
                    "duration": float(duration),
                    "fps": float(fps),
                    "resolution": f"{width}x{height}",
                    "width": width,
                    "height": height,
                    "frame_count": frame_count,
                    "format": Path(file_path).suffix.lower().lstrip('.'),
                    "codec": codec if codec.strip() else "unknown",
                    "has_audio": has_audio,
                    "file_size": len(file_content),
                    "metadata": {
                        "extracted_at": datetime.now().isoformat()
                    }
                }
            finally:
                # 清理临时文件
                try:
                    os.unlink(tmp_path)
                except:
                    pass
            
        except Exception as e:
            logger.error(f"视频解析失败: {str(e)}")
            # 如果解析失败，返回基本信息
            return {
                "file_type": "video",
                "duration": 0.0,
                "fps": 30.0,  # 默认值
                "resolution": "1920x1080",  # 默认值
                "width": 1920,
                "height": 1080,
                "frame_count": 0,
                "format": Path(file_path).suffix.lower().lstrip('.'),
                "codec": "unknown",
                "has_audio": False,
                "file_size": len(file_content),
                "metadata": {
                    "extracted_at": datetime.now().isoformat(),
                    "parse_error": str(e)
                }
            }

class ParserFactory:
    """解析器工厂"""
    
    _parsers: Dict[FileType, DocumentParser] = {
        FileType.PDF: PDFParser(),
        FileType.DOCX: DocxParser(),
        FileType.PPTX: PptxParser(),
        FileType.TXT: TextParser(),
        FileType.MD: MarkdownParser(),
        FileType.IMAGE: ImageParser(),
        FileType.AUDIO: AudioParser(),
        FileType.VIDEO: VideoParser()
    }
    
    @classmethod
    def detect_file_type(cls, file_path: str, file_content: bytes) -> FileType:
        """检测文件类型"""
        file_ext = Path(file_path).suffix.lower().lstrip('.')
        
        # 基于文件扩展名检测
        if file_ext in ["pdf"]:
            return FileType.PDF
        elif file_ext in ["docx"]:
            return FileType.DOCX
        elif file_ext in ["doc"]:
            return FileType.DOC
        elif file_ext in ["pptx"]:
            return FileType.PPTX
        elif file_ext in ["txt"]:
            return FileType.TXT
        elif file_ext in ["md", "markdown"]:
            return FileType.MD
        elif file_ext in ["jpg", "jpeg", "png", "gif", "webp", "bmp", "tiff", "tif"]:
            return FileType.IMAGE
        elif file_ext in ["mp3", "wav", "m4a", "flac", "aac", "ogg", "wma", "opus"]:
            return FileType.AUDIO
        elif file_ext in ["mp4", "avi", "mov", "mkv", "webm", "flv", "wmv", "m4v"]:
            return FileType.VIDEO
        
        # 基于内容检测（简单实现）
        try:
            # 尝试检测PDF
            if file_content.startswith(b'%PDF'):
                return FileType.PDF
            # 图片文件头：JPEG, PNG, GIF, BMP, WebP(RIFF....WEBP), TIFF(II* 或 MM\0*)
            if file_content.startswith((b'\xff\xd8\xff', b'\x89PNG', b'GIF', b'BM')):
                return FileType.IMAGE
            if len(file_content) >= 12 and file_content[:4] == b'RIFF' and file_content[8:12] == b'WEBP':
                return FileType.IMAGE
            if len(file_content) >= 4 and file_content[:2] in (b'II', b'MM') and file_content[2:4] in (b'\x2a\x00', b'\x00\x2a'):
                return FileType.IMAGE
        except Exception:
            pass
        
        return FileType.UNKNOWN
    
    @classmethod
    def get_parser(cls, file_type: FileType) -> Optional[DocumentParser]:
        """获取解析器"""
        return cls._parsers.get(file_type)
    
    @classmethod
    def register_parser(cls, file_type: FileType, parser: DocumentParser):
        """注册解析器"""
        cls._parsers[file_type] = parser
        logger.info(f"注册解析器: {file_type.value}")
    
    @classmethod
    async def parse_file(
        cls, 
        file_content: bytes, 
        file_path: str,
        **kwargs: Any,
    ) -> Dict[str, Any]:
        """解析文件；kwargs 会传给解析器（如 asset_map 供 Markdown 相对路径图片）。"""
        file_type = cls.detect_file_type(file_path, file_content)
        
        if file_type == FileType.UNKNOWN:
            raise ValueError(f"不支持的文件类型: {file_path}")
        
        parser = cls.get_parser(file_type)
        if not parser:
            raise ValueError(f"没有找到 {file_type.value} 的解析器")
        
        logger.info(f"开始解析文件: {file_path}, 类型: {file_type.value}")
        result = await parser.parse(file_content, file_path, **kwargs)
        
        # 添加文件信息
        result["file_info"] = {
            "file_path": file_path,
            "file_type": file_type.value,
            "file_size": len(file_content)
        }
        
        return result
    
    @classmethod
    def list_supported_types(cls) -> List[str]:
        """列出支持的文件类型"""
        return [ptype.value for ptype in cls._parsers.keys()]